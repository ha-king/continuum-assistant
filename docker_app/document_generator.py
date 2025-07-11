import boto3
import json
from datetime import datetime, timedelta
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import asyncio
import matplotlib.pyplot as plt
import numpy as np

class DocumentGenerator:
    def __init__(self):
        self.bedrock = boto3.client('bedrock-runtime', region_name='us-west-2')
        
    async def generate_crypto_analysis_report(self, timeframe="48-hours"):
        """Generate comprehensive crypto market analysis"""
        try:
            analysis_prompt = f"""
            Generate a comprehensive cryptocurrency market analysis and forecast for the next {timeframe}:
            
            EXECUTIVE SUMMARY:
            - Current market sentiment and key drivers
            - Major price movements in last 24 hours
            - Critical support and resistance levels
            
            TECHNICAL ANALYSIS:
            - Bitcoin (BTC) technical outlook
            - Ethereum (ETH) price action analysis  
            - Altcoin market trends (including ApeCoin)
            - Volume analysis and momentum indicators
            
            FUNDAMENTAL FACTORS:
            - Regulatory developments impacting crypto
            - Institutional adoption news
            - Macroeconomic factors (Fed policy, inflation)
            - On-chain metrics and whale activity
            
            48-HOUR FORECAST:
            - Price targets and probability scenarios
            - Key events and catalysts to watch
            - Risk factors and potential volatility triggers
            - Trading recommendations and position sizing
            
            MARKET OUTLOOK:
            - Short-term (48h) directional bias
            - Key levels to monitor
            - Contingency scenarios (bull/bear cases)
            
            Current date: {datetime.now().strftime('%Y-%m-%d %H:%M UTC')}
            
            Format with clear sections and bullet points for professional presentation.
            """
            
            response = self.bedrock.invoke_model(
                modelId="anthropic.claude-3-5-sonnet-20241022-v1:0",
                body=json.dumps({
                    "messages": [{"role": "user", "content": analysis_prompt}],
                    "max_tokens": 1500,
                    "temperature": 0.2
                })
            )
            
            result = json.loads(response['body'].read())
            return result.get('content', [{}])[0].get('text', 'Analysis generation failed')
            
        except Exception as e:
            return f"Report generation error: {str(e)}"
    
    def create_pdf_report(self, content, title="Crypto Market Analysis"):
        """Generate PDF report from content"""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch)
            
            # Styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
                textColor=colors.darkblue,
                alignment=1  # Center alignment
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                spaceAfter=12,
                textColor=colors.darkblue,
                leftIndent=0
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=8,
                leftIndent=10
            )
            
            # Build document
            story = []
            
            # Title
            story.append(Paragraph(title, title_style))
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p UTC')}", styles['Normal']))
            story.append(Spacer(1, 20))
            
            # Process content sections
            lines = content.split('\n')
            current_section = ""
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                    
                # Check if it's a heading (all caps or starts with specific keywords)
                if (line.isupper() and len(line) > 5) or line.startswith(('EXECUTIVE', 'TECHNICAL', 'FUNDAMENTAL', 'FORECAST', 'OUTLOOK')):
                    if current_section:
                        story.append(Spacer(1, 15))
                    story.append(Paragraph(line, heading_style))
                    current_section = line
                elif line.startswith('- ') or line.startswith('• '):
                    # Bullet points
                    story.append(Paragraph(line, body_style))
                else:
                    # Regular paragraphs
                    story.append(Paragraph(line, body_style))
            
            # Add footer
            story.append(Spacer(1, 30))
            footer_style = ParagraphStyle(
                'Footer',
                parent=styles['Normal'],
                fontSize=9,
                textColor=colors.grey,
                alignment=1
            )
            story.append(Paragraph("Generated by Son of Anton AI Assistant | Confidential Analysis", footer_style))
            
            doc.build(story)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            print(f"PDF generation error: {str(e)}")
            return None
    
    def create_powerpoint_presentation(self, content, title="Analysis Report"):
        """Generate PowerPoint presentation from content"""
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle_placeholder.text = f"Generated by AI Assistant\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Process content into slides
            sections = self.parse_content_sections(content)
            
            # If no clear sections, create slides from paragraphs
            if not sections or len(sections) == 1:
                sections = self.create_paragraph_sections(content)
            
            for section_title, section_content in sections.items():
                # Create slide for each section
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Title
                title_shape = slide.shapes.title
                title_shape.text = section_title[:50] + "..." if len(section_title) > 50 else section_title
                
                # Content
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                # Add content (bullet points or paragraphs)
                content_lines = [line.strip() for line in section_content if line.strip()]
                
                for i, line in enumerate(content_lines[:8]):  # Limit to 8 lines per slide
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    # Clean up bullet points
                    clean_line = line.lstrip('- •*').strip()
                    p.text = clean_line[:100] + "..." if len(clean_line) > 100 else clean_line
                    p.level = 0
                    p.font.size = Pt(16)
            
            # Save to buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            print(f"PowerPoint generation error: {str(e)}")
            return None
    
    def parse_content_sections(self, content):
        """Parse content into sections for slides"""
        sections = {}
        current_section = "Overview"
        current_content = []
        
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if it's a section header (various formats)
            is_header = (
                (line.isupper() and len(line) > 5) or
                line.startswith(('EXECUTIVE', 'TECHNICAL', 'FUNDAMENTAL', 'FORECAST', 'OUTLOOK')) or
                line.startswith('**') and line.endswith('**') or
                line.startswith('#') or
                (len(line) < 50 and ':' not in line and not line.startswith(('- ', '• ', '* ')))
            )
            
            if is_header:
                # Save previous section
                if current_content:
                    sections[current_section] = current_content
                
                # Start new section
                current_section = line.strip('*#').title()
                current_content = []
            else:
                current_content.append(line)
        
        # Save last section
        if current_content:
            sections[current_section] = current_content
        
        return sections
    
    def create_paragraph_sections(self, content):
        """Create sections from paragraphs when no clear structure exists"""
        sections = {}
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        
        for i, paragraph in enumerate(paragraphs[:10]):  # Limit to 10 slides
            # Use first few words as section title
            words = paragraph.split()[:5]
            section_title = ' '.join(words) + "..."
            sections[section_title] = [paragraph]
        
        return sections
    
    def create_chart_image(self, data_points, title="Market Trend"):
        """Create a simple chart for inclusion in documents"""
        try:
            fig, ax = plt.subplots(figsize=(8, 4))
            
            # Sample data if none provided
            if not data_points:
                x = np.linspace(0, 48, 48)  # 48 hours
                y = 45000 + np.random.normal(0, 1000, 48).cumsum()  # Sample BTC price
                data_points = list(zip(x, y))
            
            x_vals, y_vals = zip(*data_points)
            
            ax.plot(x_vals, y_vals, linewidth=2, color='#1f77b4')
            ax.set_title(title, fontsize=14, fontweight='bold')
            ax.set_xlabel('Hours')
            ax.set_ylabel('Price (USD)')
            ax.grid(True, alpha=0.3)
            
            # Save to buffer
            buffer = io.BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
            buffer.seek(0)
            plt.close()
            
            return buffer.getvalue()
            
        except Exception as e:
            print(f"Chart generation error: {str(e)}")
            return None

# Global instance
doc_generator = DocumentGenerator()

async def generate_crypto_report_pdf(timeframe="48-hours"):
    """Generate crypto analysis as PDF"""
    content = await doc_generator.generate_crypto_analysis_report(timeframe)
    pdf_data = doc_generator.create_pdf_report(content, f"Cryptocurrency Market Analysis - {timeframe.title()}")
    return pdf_data, content

async def generate_crypto_report_pptx(timeframe="48-hours"):
    """Generate crypto analysis as PowerPoint"""
    content = await doc_generator.generate_crypto_analysis_report(timeframe)
    pptx_data = doc_generator.create_powerpoint_presentation(content, f"Crypto Market Analysis - {timeframe.title()}")
    return pptx_data, content

def create_custom_report_pdf(title, content):
    """Create custom PDF report from any content"""
    return doc_generator.create_pdf_report(content, title)

def create_custom_report_pptx(title, content):
    """Create custom PowerPoint from any content"""
    return doc_generator.create_powerpoint_presentation(content, title)